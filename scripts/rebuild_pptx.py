"""Rebuild the COMP653 update slides with real chart visualizations."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = r"C:\Users\kizzi\Downloads\COMP653_Presentation_Final.pptx"
BAK = r"C:\Users\kizzi\Downloads\COMP653_Presentation_Final_BACKUP.pptx"
OUT = r"C:\Users\kizzi\Downloads\COMP653_Presentation_Updated.pptx"
FIG = r"B:\Rice\Comp653(Summer2026)\Module3\homework\stock_model\figures"

NAVY = RGBColor(0x1F, 0x28, 0x40)
BLUE = RGBColor(0x1B, 0x4F, 0x9C)
GREY = RGBColor(0x55, 0x5F, 0x70)

# Always rebuild from the pristine backup so reruns stay clean.
prs = Presentation(BAK)
SW = prs.slide_width
SH = prs.slide_height
blank = prs.slide_layouts[6]
section_layout = prs.slide_layouts[2]


def title_box(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35),
                                  SW - Inches(1.2), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    r = p.runs[0]
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = NAVY
    return tb


def caption_box(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.6), SH - Inches(1.15),
                                  SW - Inches(1.2), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    r = p.runs[0]
    r.font.size = Pt(18)
    r.font.color.rgb = GREY
    return tb


def chart_slide(title, img, caption):
    s = prs.slides.add_slide(blank)
    title_box(s, title)
    path = os.path.join(FIG, img)
    img_w = Inches(11.4)
    left = int((SW - img_w) / 2)
    s.shapes.add_picture(path, left, Inches(1.5), width=img_w)
    caption_box(s, caption)
    return s


def section_slide(title, subtitle):
    s = prs.slides.add_slide(section_layout)
    s.shapes.title.text = title
    for ph in s.placeholders:
        if ph.placeholder_format.idx != 0:
            ph.text = subtitle
            break
    return s


def bullets_slide(title, items):
    s = prs.slides.add_slide(blank)
    title_box(s, title)
    tb = s.shapes.add_textbox(Inches(0.9), Inches(1.7),
                              SW - Inches(1.8), SH - Inches(2.4))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + it
        p.space_after = Pt(14)
        r = p.runs[0]
        r.font.size = Pt(20)
        r.font.color.rgb = NAVY
    return s


new = []
new.append(section_slide("Project Update: Beyond Direction Prediction",
                         "New results produced after the original submission"))

new.append(chart_slide(
    "Direction Is Efficient, Volatility Is Not",
    "fig_direction_vs_vol.png",
    "Same pipeline, two targets. Under leak free purged validation, 90 day "
    "direction sits near the coin flip line while forward volatility reaches "
    "AUC 0.99."))

new.append(chart_slide(
    "The Signal Lives in Volatility",
    "fig_mutual_information.png",
    "Trailing volatility carries roughly seventy times the mutual information "
    "of the strongest direction feature, which is why the volatility target "
    "is learnable and direction is not."))

new.append(chart_slide(
    "The Volatility Model Is Stable, Not a Fluke",
    "fig_cpcv_paths.png",
    "Ten combinatorial purged cross validation paths, each a distinct market "
    "period, all land between 0.974 and 0.998 AUC with a mean of 0.9888."))

new.append(chart_slide(
    "One Model for the Whole Volatility Term Structure",
    "fig_term_structure.png",
    "A single model with six coupled horizon heads. Predictability rises with "
    "horizon, from 0.70 at one day to 0.98 at one month, matching how "
    "volatility smooths over longer windows."))

new.append(bullets_slide("Rigor and Engineering", [
    "Universe grew to more than 500 S&P 500 names with delisting and dollar "
    "volume liquidity filters.",
    "Training moved onto the GPU with a CuPy backend while keeping the from "
    "scratch NumPy algorithms, a large speedup on the RTX 5070.",
    "We found and fixed a label leakage gap in the final split that had "
    "inflated earlier held out numbers, then re measured everything.",
    "Every reported number is leak free under purge and embargo.",
]))

# Move the new slides to sit right after slide index 14.
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
n_new = len(new)
new_ids = ids[-n_new:]
for el in new_ids:
    sldIdLst.remove(el)
for offset, el in enumerate(new_ids):
    sldIdLst.insert(15 + offset, el)

# Update the outdated results and takeaways speaker notes.
n11 = prs.slides[11].notes_slide.notes_text_frame
if "UPDATE:" not in n11.text:
    n11.text = n11.text + (
        "\n\nUPDATE: These direction numbers were the original submission. "
        "Later combinatorial purged cross validation showed 90 day direction "
        "sits near 0.50 AUC out of sample. We reframed the target to forward "
        "realized volatility and reached AUC 0.99, which is the result we now "
        "lead with.")
n14 = prs.slides[14].notes_slide.notes_text_frame
if "UPDATE ON NEXT STEPS:" not in n14.text:
    n14.text = n14.text + (
        "\n\nUPDATE ON NEXT STEPS: The LSTM branch, market neutral sector "
        "clustering, longer training, and 5 step PGD are now implemented. The "
        "decisive move was reframing to volatility, which reached AUC 0.99, "
        "and a single coupled model for the whole volatility term structure.")

prs.save(OUT)
print("Rebuilt", OUT, "with", len(prs.slides), "slides")
